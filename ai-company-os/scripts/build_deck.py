#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_deck.py — デッキ仕様(JSON)から .pptx を生成する汎用ビルダー

入力: デッキ仕様 JSON(templates/deck_spec_example.json 参照)
出力: docs/04_PowerPoint.md 標準に準拠した 16:9 の .pptx

使い方:
  python scripts/build_deck.py spec.json output.pptx
  python scripts/build_deck.py spec.json output.pptx && python scripts/verify_pptx.py output.pptx

対応レイアウト:
  title       — 表紙(ダーク背景、見出し・サブコピー・脚注)
  bullets     — 見出し+ブレット(最大5、出典行つき)。多用しすぎない(下記ルール参照)
  cards       — カードグリッド(数字・ラベル・本文、2〜4枚)
  steps       — 番号バッジつき横並びステップ(3〜4個、カード間に矢印つき)
  agenda      — アジェンダ・章区切り(ダーク背景、番号付きリスト)
  comparison  — 2カラム比較(自社/他社、Before/After など)
  table       — 表(価格表・比較マトリクス等、ヘッダー行+データ行)
  quote       — 引用・お客様の声(大見出し引用+出典)
  bar_chart   — ネイティブ棒グラフ(編集可能、出典行つき)
  line_chart  — ネイティブ折れ線グラフ(推移データ、出典行つき)
  pie_chart   — ネイティブ円グラフ(構成比、出典行つき)
  cta         — クロージング(ダーク背景、CTA 1つ)
  big_stat    — 1スライド1つの数字だけに焦点を当てる、視覚重視レイアウト(ダーク背景)
  diagram     — ノード(丸/角丸ボックス)を矢印でつなぐ、簡易な概念図レイアウト(2〜4ノード)

パレットプリセット(meta.palette_preset で指定、meta.palette で個別上書き可):
  navy_gold    — 既定。紺+氷色+金(信頼・フォーマル)
  forest_sand  — 深緑+砂色+テラコッタ(サステナ・自然派)
  slate_azure  — スレート+空色+アズール(テック・プロダクト)

規則(ビルド時に強制):
  - 全スライドに note(スピーカーノート)必須。欠けていればエラー終了
  - bullets は最大5個。超えていればエラー終了
  - table の各行の列数はヘッダーと一致していること。不一致ならエラー終了
  - bullets レイアウトはデッキ全体の40%まで(10枚以上のデッキで適用)。超えていれば
    エラー終了し、cards/diagram/big_stat/steps等への置き換えを促す(文字量を減らし
    図解・グラフで見て分かるようにするため。docs/04_PowerPoint.md参照)
  - diagram の nodes は2つ以上、big_stat の stat は空でないこと
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx が見つかりません。`pip install python-pptx` を実行してください。")
    sys.exit(2)

W, H = 13.333, 7.5

PALETTE_PRESETS = {
    "navy_gold": {
        "primary": "1E2761",
        "primary_dark": "141B47",
        "base": "CADCFC",
        "base_soft": "8FA3D0",
        "accent": "D9A441",
        "ink": "1A1A2E",
        "muted": "5B6270",
        "card": "F3F6FC",
    },
    "forest_sand": {
        "primary": "24523F",
        "primary_dark": "17362A",
        "base": "E7DCC2",
        "base_soft": "BFAE84",
        "accent": "C1622D",
        "ink": "20261F",
        "muted": "5E6B5C",
        "card": "F4F1E6",
    },
    "slate_azure": {
        "primary": "2B3A55",
        "primary_dark": "1B2537",
        "base": "D3E4F5",
        "base_soft": "9CB8D6",
        "accent": "1F8FD6",
        "ink": "1C2230",
        "muted": "5A6474",
        "card": "F1F5FA",
    },
}
DEFAULT_PALETTE = PALETTE_PRESETS["navy_gold"]
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def rgb(hexstr: str) -> RGBColor:
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


SHADOW_GRAY = RGBColor(0xDD, 0xDD, 0xE3)


class DeckBuilder:
    def __init__(self, meta: dict, total_slides: int = 0):
        preset = PALETTE_PRESETS.get(meta.get("palette_preset", "navy_gold"), DEFAULT_PALETTE)
        pal = {**preset, **meta.get("palette", {})}
        self.c = {k: rgb(v) for k, v in pal.items()}
        self.font = meta.get("font", "Yu Gothic")
        self.footer_text = meta.get("footer", "")
        self.total_slides = total_slides
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.blank = self.prs.slide_layouts[6]
        self.page = 0

    # ---- 低レベルヘルパー ----

    def slide(self, dark: bool = False):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = self.c["primary_dark"] if dark else WHITE
        self.page += 1
        self.progress_bar(s)
        return s

    def progress_bar(self, s):
        """スライド最上部に、全体のうち今どこまで進んだかを示す細い進捗バーを描く。
        「あとどれくらいで終わるか」が視覚的にわかることで最後まで見る動機を作る
        (docs/04_PowerPoint.md エンゲージメント設計ルール参照)。"""
        if not self.total_slides:
            return
        bar_h = 0.05
        track = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(bar_h))
        track.fill.solid()
        track.fill.fore_color.rgb = self.c["base_soft"]
        track.line.fill.background()
        track.shadow.inherit = False
        filled_w = max(W * (self.page / self.total_slides), 0.001)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(filled_w), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.c["accent"]
        bar.line.fill.background()
        bar.shadow.inherit = False

    def text(self, s, txt, x, y, w, h, size=14, color=None, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(str(txt).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if spacing != 1.0:
                p.line_spacing = spacing
            r = p.add_run()
            r.text = line
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.name = self.font
            f.color.rgb = color if color is not None else self.c["ink"]
        return tb

    def card(self, s, x, y, w, h):
        # 立体感を出すため、本体の少し右下にずらした影用の矩形を先に描いてから本体を重ねる
        shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(y + 0.06),
                                     Inches(w), Inches(h))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SHADOW_GRAY
        shadow.line.fill.background()
        shadow.adjustments[0] = 0.06
        shadow.shadow.inherit = False

        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self.c["card"]
        sh.line.fill.background()
        sh.adjustments[0] = 0.06
        sh.shadow.inherit = False
        return sh

    def badge(self, s, glyph, x, y, d=0.7, size=20):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        c.fill.solid()
        c.fill.fore_color.rgb = self.c["primary"]
        c.line.fill.background()
        tf = c.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(glyph)
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.name = self.font
        r.font.color.rgb = WHITE
        return c

    def title_bar(self, s, t):
        self.text(s, t, 0.5, 0.45, W - 1.0, 0.8, size=30, bold=True, color=self.c["primary"])
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(0.55), Inches(0.06))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.c["accent"]
        accent.line.fill.background()
        accent.shadow.inherit = False

    def footer(self, s):
        if self.footer_text:
            self.text(s, self.footer_text, 0.5, 7.05, 8, 0.35, size=10, color=self.c["muted"])
        page_label = f"{self.page} / {self.total_slides}" if self.total_slides else str(self.page)
        self.text(s, page_label, W - 1.5, 7.05, 1.0, 0.35, size=10,
                  color=self.c["muted"], align=PP_ALIGN.RIGHT)

    def source_line(self, s, src):
        if src:
            self.text(s, f"出典: {src}", 0.5, 6.6, W - 1.0, 0.35, size=10, color=self.c["muted"])

    def notes(self, s, t):
        s.notes_slide.notes_text_frame.text = t

    # ---- レイアウト ----

    def add_title(self, spec):
        s = self.slide(dark=True)
        self.text(s, spec["heading"], 0.6, 2.2, 12.1, 2.3, size=42, bold=True, color=WHITE, spacing=1.15)
        if spec.get("sub"):
            self.text(s, spec["sub"], 0.6, 4.6, 11.5, 0.6, size=20, color=self.c["base"])
        if spec.get("footnote"):
            self.text(s, spec["footnote"], 0.6, 6.5, 11.5, 0.4, size=12, color=self.c["base_soft"])
        return s

    def add_bullets(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        y = 1.6
        if spec.get("lead"):
            self.text(s, spec["lead"], 0.5, y, W - 1.0, 0.5, size=16, color=self.c["muted"])
            y += 0.7
        for b in spec["bullets"]:
            self.badge(s, "・", 0.5, y + 0.02, d=0.35, size=12)
            self.text(s, b, 1.05, y, W - 1.6, 0.7, size=18, spacing=1.1)
            y += 0.85
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_cards(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        cards = spec["cards"]
        n = len(cards)
        gap = 0.35
        cw = (W - 1.0 - gap * (n - 1)) / n
        for i, cd in enumerate(cards):
            x = 0.5 + i * (cw + gap)
            self.card(s, x, 1.7, cw, 4.3)
            y = 2.0
            if cd.get("big"):
                self.text(s, cd["big"], x + 0.3, y, cw - 0.6, 0.9, size=34, bold=True, color=self.c["accent"])
                y += 1.1
            self.text(s, cd["label"], x + 0.3, y, cw - 0.6, 0.6, size=16, bold=True, color=self.c["primary"])
            y += 0.75
            if cd.get("text"):
                self.text(s, cd["text"], x + 0.3, y, cw - 0.6, 2.0, size=13, color=self.c["muted"], spacing=1.15)
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_steps(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        steps = spec["steps"]
        n = len(steps)
        gap = 0.4
        cw = (W - 1.0 - gap * (n - 1)) / n
        for i, st in enumerate(steps):
            x = 0.5 + i * (cw + gap)
            self.card(s, x, 1.9, cw, 3.6)
            self.badge(s, i + 1, x + 0.3, 2.2, d=0.7, size=20)
            self.text(s, st["label"], x + 0.3, 3.2, cw - 0.6, 0.6, size=16, bold=True, color=self.c["primary"])
            self.text(s, st.get("text", ""), x + 0.3, 3.9, cw - 0.6, 1.4, size=12, color=self.c["muted"], spacing=1.15)
            if i < n - 1:
                arrow_x = x + cw + gap * 0.18
                arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(arrow_x), Inches(3.5), Inches(gap * 0.64), Inches(0.4))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.c["accent"]
                arrow.line.fill.background()
        self.footer(s)
        return s

    def add_big_stat(self, spec):
        """1スライド1つの数字だけに焦点を当てる、視覚重視のレイアウト。"""
        s = self.slide(dark=True)
        self.text(s, spec["title"], 0.6, 0.6, 12.1, 0.7, size=22, bold=True, color=self.c["base"])
        self.text(s, spec["stat"], 0.6, 2.0, 12.1, 2.2, size=88, bold=True, color=self.c["accent"],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if spec.get("caption"):
            self.text(s, spec["caption"], 1.2, 4.4, 10.9, 0.8, size=20, color=WHITE,
                      align=PP_ALIGN.CENTER, spacing=1.2)
        if spec.get("source"):
            self.text(s, f"出典: {spec['source']}", 0.6, 6.6, 12.1, 0.4, size=11,
                      color=self.c["base_soft"], align=PP_ALIGN.CENTER)
        return s

    def add_diagram(self, spec):
        """ノード(丸/角丸ボックス)を矢印でつなぐ、簡易な概念図レイアウト。"""
        s = self.slide()
        self.title_bar(s, spec["title"])
        nodes = spec["nodes"]
        n = len(nodes)
        gap = 0.9
        nw = 2.6
        total_w = nw * n + gap * (n - 1)
        start_x = (W - total_w) / 2
        y = 2.8
        nh = 1.6
        for i, node in enumerate(nodes):
            x = start_x + i * (nw + gap)
            self.card(s, x, y, nw, nh)
            self.text(s, node["label"], x + 0.2, y + 0.2, nw - 0.4, 0.7, size=17, bold=True,
                      color=self.c["primary"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if node.get("sublabel"):
                self.text(s, node["sublabel"], x + 0.2, y + 0.9, nw - 0.4, 0.6, size=12,
                          color=self.c["muted"], align=PP_ALIGN.CENTER, spacing=1.1)
            if i < n - 1:
                ax = x + nw + 0.08
                arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(y + nh / 2 - 0.2),
                                            Inches(gap - 0.16), Inches(0.4))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.c["accent"]
                arrow.line.fill.background()
        if spec.get("caption"):
            self.text(s, spec["caption"], 0.6, y + nh + 0.5, W - 1.2, 0.8, size=14,
                      color=self.c["muted"], align=PP_ALIGN.CENTER, spacing=1.2)
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_bar_chart(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        data = CategoryChartData()
        data.categories = spec["categories"]
        data.add_series(spec["series"].get("name", ""), spec["series"]["values"])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.7), data)
        chart = gf.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = spec.get("number_format", "#,##0")
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.size = Pt(12)
        dl.font.bold = True
        dl.font.name = self.font
        dl.font.color.rgb = self.c["primary"]
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = self.c["primary"]
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_cta(self, spec):
        s = self.slide(dark=True)
        self.text(s, spec["heading"], 0.6, 1.8, 12.1, 1.6, size=36, bold=True, color=WHITE, spacing=1.15)
        if spec.get("sub"):
            self.text(s, spec["sub"], 0.6, 3.6, 11.5, 1.0, size=18, color=self.c["base"], spacing=1.2)
        if spec.get("action"):
            self.card(s, 0.6, 5.0, 8.0, 1.1)
            self.text(s, spec["action"], 0.95, 5.28, 7.4, 0.6, size=18, bold=True, color=self.c["primary"])
        return s

    def add_agenda(self, spec):
        s = self.slide(dark=True)
        self.text(s, spec["title"], 0.6, 0.7, 12.1, 1.0, size=32, bold=True, color=WHITE)
        y = 2.1
        for i, item in enumerate(spec["items"], start=1):
            self.badge(s, i, 0.6, y, d=0.6, size=18)
            self.text(s, item, 1.45, y + 0.02, 10.5, 0.6, size=20, color=self.c["base"], anchor=MSO_ANCHOR.MIDDLE)
            y += 0.85
        return s

    def add_comparison(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        cols = [spec["left"], spec["right"]]
        gap = 0.4
        cw = (W - 1.0 - gap) / 2
        for i, col in enumerate(cols):
            x = 0.5 + i * (cw + gap)
            self.card(s, x, 1.7, cw, 4.6)
            self.text(s, col["heading"], x + 0.35, 2.0, cw - 0.7, 0.6, size=18, bold=True, color=self.c["primary"])
            y = 2.75
            for item in col.get("items", []):
                self.badge(s, "・", x + 0.35, y + 0.02, d=0.3, size=11)
                self.text(s, item, x + 0.8, y, cw - 1.15, 0.6, size=14, spacing=1.1)
                y += 0.6
        self.footer(s)
        return s

    def add_table(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        headers = spec["headers"]
        rows = spec["rows"]
        n_rows, n_cols = len(rows) + 1, len(headers)
        top, height = 1.7, min(4.5, 0.55 * (n_rows))
        gf = s.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(top), Inches(W - 1.0), Inches(height))
        table = gf.table
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.c["primary"]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.name = self.font
            para.font.color.rgb = WHITE
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.c["card"] if r % 2 == 0 else WHITE
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.font.name = self.font
                para.font.color.rgb = self.c["ink"]
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_quote(self, spec):
        s = self.slide(dark=True)
        self.text(s, f"“{spec['quote']}”", 1.2, 2.0, 10.9, 2.8, size=28, bold=True,
                  color=WHITE, spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
        if spec.get("attribution"):
            self.text(s, f"— {spec['attribution']}", 1.2, 5.0, 10.9, 0.6, size=16, color=self.c["base"])
        return s

    def add_pie_chart(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        data = CategoryChartData()
        data.categories = spec["categories"]
        data.add_series(spec.get("series_name", ""), spec["values"])
        gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(2.4), Inches(1.6), Inches(8.5), Inches(4.7), data)
        chart = gf.chart
        chart.has_legend = True
        chart.legend.position = 2  # XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = self.font
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = spec.get("number_format", "0%")
        dl.number_format_is_linked = False
        dl.font.size = Pt(11)
        dl.font.bold = True
        dl.font.name = self.font
        dl.font.color.rgb = WHITE
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s

    def add_line_chart(self, spec):
        s = self.slide()
        self.title_bar(s, spec["title"])
        data = CategoryChartData()
        data.categories = spec["categories"]
        data.add_series(spec["series"].get("name", ""), spec["series"]["values"])
        gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.7), data)
        chart = gf.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = spec.get("number_format", "#,##0")
        dl.number_format_is_linked = False
        dl.font.size = Pt(11)
        dl.font.name = self.font
        dl.font.color.rgb = self.c["primary"]
        series = chart.series[0]
        series.format.line.color.rgb = self.c["accent"]
        series.format.line.width = Pt(2.5)
        self.source_line(s, spec.get("source"))
        self.footer(s)
        return s


LAYOUTS = {
    "title": DeckBuilder.add_title,
    "bullets": DeckBuilder.add_bullets,
    "cards": DeckBuilder.add_cards,
    "steps": DeckBuilder.add_steps,
    "agenda": DeckBuilder.add_agenda,
    "comparison": DeckBuilder.add_comparison,
    "table": DeckBuilder.add_table,
    "quote": DeckBuilder.add_quote,
    "bar_chart": DeckBuilder.add_bar_chart,
    "line_chart": DeckBuilder.add_line_chart,
    "pie_chart": DeckBuilder.add_pie_chart,
    "cta": DeckBuilder.add_cta,
    "big_stat": DeckBuilder.add_big_stat,
    "diagram": DeckBuilder.add_diagram,
}

MAX_BULLETS = 5
MAX_CONSECUTIVE_SAME_LAYOUT = 3  # これを超えて同一レイアウトが連続すると視聴離脱を招きやすい(docs/04_PowerPoint.md)
MAX_BULLETS_RATIO = 0.4  # bullets(文字中心)レイアウトはデッキ全体の40%まで(docs/04_PowerPoint.md)
MIN_SLIDES_FOR_RATIO_CHECK = 10  # 短いデッキ(フィクスチャ等)ではこの比率チェックを適用しない


def build(spec_path: str, out_path: str) -> int:
    spec = json.loads(Path(spec_path).read_text(encoding="utf-8"))
    errors: list[str] = []
    for i, sl in enumerate(spec.get("slides", []), start=1):
        if sl.get("layout") not in LAYOUTS:
            errors.append(f"スライド {i}: 未知のレイアウト '{sl.get('layout')}'")
        if not str(sl.get("note", "")).strip():
            errors.append(f"スライド {i}: note(スピーカーノート)がありません")
        if sl.get("layout") == "bullets" and len(sl.get("bullets", [])) > MAX_BULLETS:
            errors.append(f"スライド {i}: ブレットが {MAX_BULLETS} 個を超えています")
        if sl.get("layout") == "table":
            n_headers = len(sl.get("headers", []))
            for r, row in enumerate(sl.get("rows", []), start=1):
                if len(row) != n_headers:
                    errors.append(f"スライド {i}: table の行 {r} の列数がヘッダー({n_headers}列)と一致しません")
        if sl.get("layout") == "diagram" and len(sl.get("nodes", [])) < 2:
            errors.append(f"スライド {i}: diagram の nodes は2つ以上必要です")
        if sl.get("layout") == "big_stat" and not str(sl.get("stat", "")).strip():
            errors.append(f"スライド {i}: big_stat の stat が空です")
    if not spec.get("slides"):
        errors.append("slides が空です")

    run_layout, run_len, run_start = None, 0, 0
    for i, sl in enumerate(spec.get("slides", []), start=1):
        layout = sl.get("layout")
        if layout == run_layout:
            run_len += 1
        else:
            run_layout, run_len, run_start = layout, 1, i
        if run_len == MAX_CONSECUTIVE_SAME_LAYOUT + 1:
            errors.append(
                f"スライド {run_start}〜{i}: 同一レイアウト '{layout}' が"
                f"{MAX_CONSECUTIVE_SAME_LAYOUT}枚を超えて連続しています"
                "(視聴離脱防止のためレイアウトに変化をつけること。docs/04_PowerPoint.md参照)"
            )

    all_slides = spec.get("slides", [])
    if len(all_slides) >= MIN_SLIDES_FOR_RATIO_CHECK:
        n_bullets = sum(1 for sl in all_slides if sl.get("layout") == "bullets")
        ratio = n_bullets / len(all_slides)
        if ratio > MAX_BULLETS_RATIO:
            errors.append(
                f"bullets(文字中心)レイアウトが全{len(all_slides)}枚中{n_bullets}枚"
                f"({ratio:.0%})で、上限の{MAX_BULLETS_RATIO:.0%}を超えています。"
                "cards/diagram/big_stat/steps/comparison等の視覚的なレイアウトに"
                "置き換えること(docs/04_PowerPoint.md参照)"
            )

    if errors:
        print(f"仕様エラー ({len(errors)} 件):")
        for e in errors:
            print(f"  - {e}")
        return 1

    builder = DeckBuilder(spec.get("meta", {}), total_slides=len(spec["slides"]))
    for sl in spec["slides"]:
        s = LAYOUTS[sl["layout"]](builder, sl)
        builder.notes(s, sl["note"])
    builder.prs.save(out_path)
    print(f"生成完了: {out_path} ({len(spec['slides'])} スライド)")
    print(f"次に検品を実行してください: python scripts/verify_pptx.py {out_path}")
    return 0


def main() -> int:
    if len(sys.argv) != 3:
        print("使い方: python scripts/build_deck.py <spec.json> <output.pptx>")
        return 2
    if not Path(sys.argv[1]).is_file():
        print(f"ERROR: 仕様ファイルが存在しません: {sys.argv[1]}")
        return 2
    return build(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    sys.exit(main())
