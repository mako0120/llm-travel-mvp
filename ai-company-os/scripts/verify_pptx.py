#!/usr/bin/env python3
"""verify_pptx.py — .pptx 成果物の自動検品

入力: .pptx ファイルパス / 出力: 検品レポート(stdout)と終了コード(0=合格, 1=不合格, 2=実行エラー)

検品項目 (docs/04_PowerPoint.md):
  1. ファイルが python-pptx で開ける(破損していない)
  2. スライド枚数が指定範囲内 (--min-slides / --max-slides)
  3. 全スライドにスピーカーノートがある (--require-notes, 既定で有効)
  4. 図形がスライド境界からはみ出していない
  5. 空のプレースホルダが残っていない

使い方:
  python scripts/verify_pptx.py deck.pptx
  python scripts/verify_pptx.py deck.pptx --min-slides 11 --max-slides 11
  python scripts/verify_pptx.py --self-test
"""

from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: python-pptx が見つかりません。`pip install python-pptx` を実行してください。")
    sys.exit(2)


def check_deck(path: str, min_slides: int, max_slides: int | None, require_notes: bool) -> list[str]:
    """検品を実行し、問題点のリストを返す(空リスト=合格)。"""
    issues: list[str] = []

    try:
        prs = Presentation(path)
    except Exception as exc:  # 破損・非pptx
        return [f"ファイルを開けません: {exc}"]

    slides = list(prs.slides)
    n = len(slides)
    if n < min_slides:
        issues.append(f"スライド枚数 {n} 枚は最小 {min_slides} 枚を下回っています")
    if max_slides is not None and n > max_slides:
        issues.append(f"スライド枚数 {n} 枚は最大 {max_slides} 枚を超えています")

    slide_w, slide_h = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(slides, start=1):
        if require_notes:
            notes = ""
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if not notes:
                issues.append(f"スライド {idx}: スピーカーノートがありません")

        for shape in slide.shapes:
            name = shape.name or shape.shape_type
            if None in (shape.left, shape.top, shape.width, shape.height):
                continue  # 位置未定義の図形(コネクタ等)は境界判定不能
            if (
                shape.left < Emu(0)
                or shape.top < Emu(0)
                or shape.left + shape.width > slide_w
                or shape.top + shape.height > slide_h
            ):
                issues.append(f"スライド {idx}: 図形 '{name}' がスライド境界からはみ出しています")

            if shape.is_placeholder and shape.has_text_frame:
                if not shape.text_frame.text.strip():
                    issues.append(f"スライド {idx}: 空のプレースホルダ '{name}' が残っています")

    return issues


def self_test() -> int:
    """合格すべきデッキと不合格すべきデッキを生成し、判定が正しいことを確認する。"""
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as tmp:
        good = str(Path(tmp) / "good.pptx")
        bad = str(Path(tmp) / "bad.pptx")

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = "Self test slide"
        slide.notes_slide.notes_text_frame.text = "テスト用ノート"
        prs.save(good)

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(5), Inches(2))  # はみ出し
        box.text_frame.text = "overflow"
        prs.save(bad)  # ノートなし + はみ出し

        good_issues = check_deck(good, min_slides=1, max_slides=None, require_notes=True)
        bad_issues = check_deck(bad, min_slides=2, max_slides=None, require_notes=True)

        ok = True
        if good_issues:
            print(f"SELF-TEST FAIL: 合格すべきデッキで問題検出: {good_issues}")
            ok = False
        expected = ["スライド枚数", "スピーカーノート", "はみ出し"]
        for key in expected:
            if not any(key in i for i in bad_issues):
                print(f"SELF-TEST FAIL: 不合格デッキで '{key}' が検出されませんでした: {bad_issues}")
                ok = False

        print("SELF-TEST PASSED" if ok else "SELF-TEST FAILED")
        return 0 if ok else 1


def main() -> int:
    parser = argparse.ArgumentParser(description=".pptx 成果物の自動検品")
    parser.add_argument("pptx", nargs="?", help="検品する .pptx ファイル")
    parser.add_argument("--min-slides", type=int, default=1)
    parser.add_argument("--max-slides", type=int, default=None)
    parser.add_argument("--no-require-notes", action="store_true", help="ノート必須チェックを外す")
    parser.add_argument("--self-test", action="store_true", help="検品ロジック自体を検証する")
    args = parser.parse_args()

    if args.self_test:
        return self_test()

    if not args.pptx:
        parser.error("検品する .pptx ファイルか --self-test を指定してください")
    if not Path(args.pptx).is_file():
        print(f"ERROR: ファイルが存在しません: {args.pptx}")
        return 2

    issues = check_deck(
        args.pptx,
        min_slides=args.min_slides,
        max_slides=args.max_slides,
        require_notes=not args.no_require_notes,
    )

    print(f"検品対象: {args.pptx}")
    if issues:
        print(f"不合格 ({len(issues)} 件):")
        for issue in issues:
            print(f"  - {issue}")
        return 1
    print("合格: 問題は検出されませんでした")
    return 0


if __name__ == "__main__":
    sys.exit(main())
